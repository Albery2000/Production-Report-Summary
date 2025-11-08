import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import io
import base64

def extract_wells_with_net_diff_bo(file_content):
    """
    Extract wells that have Net Diff BO values (excluding zeros) from specific columns and stop at TOTAL row
    """
    try:
        # Read the Excel file with multi-level headers, skipping first 6 rows
        df = pd.read_excel(
            file_content, 
            sheet_name='Report', 
            skiprows=6,
            header=[0, 1]  # Two header rows
        )
        
        st.subheader("🔍 Detected Column Structure")
        
        # Display all columns to help with debugging
        columns_info = []
        for i, col in enumerate(df.columns):
            col_info = {
                'Column Index': i,
                'Level 0': str(col[0]) if pd.notna(col[0]) else '',
                'Level 1': str(col[1]) if len(col) > 1 and pd.notna(col[1]) else '',
                'Full Name': str(col)
            }
            columns_info.append(col_info)
        
        columns_df = pd.DataFrame(columns_info)
        st.dataframe(columns_df)
        
        # Find the specific columns we need
        field_col = None
        well_name_col = None
        net_diff_bo_col = None
        net_bo_col = None
        
        for i, col in enumerate(df.columns):
            # Look for the exact column structures
            if str(col) == "('TOTAL PRODUCTION', 'Net diff. BO')":
                net_diff_bo_col = col
                st.success(f"✅ Found Net Diff BO column: {col} (Index {i})")
            
            elif str(col) == "('TOTAL PRODUCTION', 'Net\\nBO')" or "('TOTAL PRODUCTION', 'Net\nBO')" in str(col):
                net_bo_col = col
                st.success(f"✅ Found Net BO column: {col} (Index {i})")
            
            # Field column - look for ('Field', 'Unnamed: 0_level_1')
            elif str(col) == "('Field', 'Unnamed: 0_level_1')":
                field_col = col
                st.success(f"✅ Found Field column: {col} (Index {i})")
            
            # Well name column - look for ('RUNNING WELLS', 'Unnamed: 1_level_1')
            elif str(col) == "('RUNNING WELLS', 'Unnamed: 1_level_1')":
                well_name_col = col
                st.success(f"✅ Found Well Name column: {col} (Index {i})")
        
        # Validation
        if field_col is None:
            st.error("❌ Could not find 'Field' column")
            return None, None, None, None, None
        
        if well_name_col is None:
            st.error("❌ Could not find well name column")
            return None, None, None, None, None
        
        if net_diff_bo_col is None:
            st.error("❌ Could not find 'Net diff. BO' column")
            return None, None, None, None, None
        
        if net_bo_col is None:
            st.error("❌ Could not find 'Net BO' column")
            return None, None, None, None, None
        
        # Convert numeric columns
        df[net_diff_bo_col] = pd.to_numeric(df[net_diff_bo_col], errors='coerce')
        df[net_bo_col] = pd.to_numeric(df[net_bo_col], errors='coerce')
        
        # Find where to stop (at "TOTAL" in Field column)
        stop_index = None
        for idx, value in enumerate(df[field_col]):
            if pd.notna(value) and 'TOTAL' in str(value).upper():
                stop_index = idx
                st.info(f"🛑 Found 'TOTAL' row at index {idx}, stopping extraction here")
                break
        
        # If no TOTAL found, use all rows
        if stop_index is None:
            stop_index = len(df)
            st.warning("⚠️ No 'TOTAL' row found, using all available data")
        
        # Filter rows up to the TOTAL row
        df_before_total = df.iloc[:stop_index].copy()
        
        # Calculate TOTAL statistics for ALL wells (including zeros)
        all_wells_count = len(df_before_total)
        total_net_bo_all = df_before_total[net_bo_col].sum()
        total_net_diff_bo_all = df_before_total[net_diff_bo_col].sum()
        
        # Filter rows that have Net diff. BO values AND are not zero (but include negative values)
        filtered_df = df_before_total[
            (df_before_total[net_diff_bo_col].notna()) & 
            (df_before_total[net_diff_bo_col] != 0)  # Exclude zeros but include negatives
        ].copy()
        
        if filtered_df.empty:
            st.warning("⚠️ No wells found with non-zero Net Diff BO values before TOTAL row")
            return None, None, None, None, None
        
        # Show how many wells were filtered out due to zero values
        all_wells_with_net_diff = df_before_total[df_before_total[net_diff_bo_col].notna()]
        zero_wells_count = len(all_wells_with_net_diff[all_wells_with_net_diff[net_diff_bo_col] == 0])
        st.info(f"📊 Filtered out {zero_wells_count} wells with zero Net Diff BO values")
        
        # Show distribution of positive vs negative values
        positive_count = len(filtered_df[filtered_df[net_diff_bo_col] > 0])
        negative_count = len(filtered_df[filtered_df[net_diff_bo_col] < 0])
        st.info(f"📈 Value distribution: {positive_count} positive, {negative_count} negative Net Diff BO values")
        
        # Select the columns we need in the correct order
        result_columns = [field_col, well_name_col, net_bo_col, net_diff_bo_col]
        
        # Create final result dataframe
        result_df = filtered_df[result_columns].copy()
        
        # Clean up the data - remove rows where well name is empty or is a field name
        field_names = ['Ferdaus', 'Sidra', 'Ganna', 'Rayan', 'Abrar', 'Abrar-South', 'Rawda']
        
        # Filter out rows where well_name is actually a field name
        mask = ~result_df[well_name_col].isin(field_names)
        result_df = result_df[mask].copy()
        
        # Remove rows where well_name is empty or NaN
        result_df = result_df[result_df[well_name_col].notna()]
        result_df = result_df[result_df[well_name_col] != '']
        
        result_df = result_df.reset_index(drop=True)
        
        # Calculate totals and statistics for non-zero wells
        total_net_bo_non_zero = result_df[net_bo_col].sum()
        total_net_diff_bo_non_zero = result_df[net_diff_bo_col].sum()
        well_count_non_zero = len(result_df)
        
        # Calculate statistics for both ALL wells and non-zero wells
        stats = {
            # All Wells Statistics
            'Total All Wells': all_wells_count,
            'Total Net BO (All Wells)': total_net_bo_all,
            'Total Net Diff BO (All Wells)': total_net_diff_bo_all,
            'Average Net BO (All Wells)': df_before_total[net_bo_col].mean(),
            'Average Net Diff BO (All Wells)': df_before_total[net_diff_bo_col].mean(),
            
            # Non-Zero Wells Statistics
            'Total Wells with Non-Zero Net Diff BO': well_count_non_zero,
            'Positive Net Diff BO Wells': positive_count,
            'Negative Net Diff BO Wells': negative_count,
            'Total Net BO (Non-Zero Wells)': total_net_bo_non_zero,
            'Total Net Diff BO (Non-Zero Wells)': total_net_diff_bo_non_zero,
            'Average Net BO (Non-Zero Wells)': result_df[net_bo_col].mean(),
            'Average Net Diff BO (Non-Zero Wells)': result_df[net_diff_bo_col].mean(),
            'Maximum Net BO': result_df[net_bo_col].max(),
            'Maximum Net Diff BO': result_df[net_diff_bo_col].max(),
            'Minimum Net BO': result_df[net_bo_col].min(),
            'Minimum Net Diff BO': result_df[net_diff_bo_col].min(),
            'Median Net BO': result_df[net_bo_col].median(),
            'Median Net Diff BO': result_df[net_diff_bo_col].median(),
            'Standard Deviation Net BO': result_df[net_bo_col].std(),
            'Standard Deviation Net Diff BO': result_df[net_diff_bo_col].std()
        }
        
        # Create the final dataframe with proper column structure
        final_df = result_df.copy()
        
        # Format numeric columns
        for col in [net_bo_col, net_diff_bo_col]:
            if col in final_df.columns and final_df[col].dtype in [np.float64, np.int64]:
                final_df[col] = final_df[col].round(2)
        
        # Add TOTAL (All Wells) row with net bo and net diff bo
        total_row_all = pd.DataFrame([{
            field_col: 'TOTAL (All Wells)',
            well_name_col: f'{all_wells_count} Total Wells',
            net_bo_col: total_net_bo_all,
            net_diff_bo_col: total_net_diff_bo_all
        }])
        
        # Combine main data with total row
        final_df = pd.concat([final_df, total_row_all], ignore_index=True)
        
        st.success(f"✅ Successfully extracted {well_count_non_zero} wells with non-zero Net Diff BO values")
        
        return final_df, well_count_non_zero, stats, [field_col, well_name_col, net_bo_col, net_diff_bo_col], df_before_total
        
    except Exception as e:
        st.error(f"❌ Error processing file: {str(e)}")
        import traceback
        st.error(f"Detailed error: {traceback.format_exc()}")
        return None, None, None, None, None

def create_visualizations(data_without_total, original_columns, all_wells_data):
    """
    Create simplified statistical visualizations with only three charts
    """
    try:
        # Check if we have valid data for visualizations
        if data_without_total.empty or all_wells_data.empty:
            st.warning("No data available for visualizations")
            return None
            
        # Extract the column names
        field_col = original_columns[0]      # ('Field', 'Unnamed: 0_level_1')
        well_name_col = original_columns[1]  # ('RUNNING WELLS', 'Unnamed: 1_level_1')
        net_bo_col = original_columns[2]     # ('TOTAL PRODUCTION', 'Net\nBO')
        net_diff_bo_col = original_columns[3] # ('TOTAL PRODUCTION', 'Net diff. BO')
        
        # Create clean copies for visualization
        viz_data_non_zero = data_without_total.copy()
        viz_data_all = all_wells_data.copy()
        
        # Remove rows with NaN values in the key columns for visualization
        viz_data_non_zero = viz_data_non_zero[
            viz_data_non_zero[well_name_col].notna() & 
            viz_data_non_zero[net_bo_col].notna() & 
            viz_data_non_zero[net_diff_bo_col].notna()
        ]
        
        viz_data_all = viz_data_all[
            viz_data_all[well_name_col].notna() & 
            viz_data_all[net_bo_col].notna()
        ]
        
        # Check if we have any data left after cleaning
        if viz_data_non_zero.empty or viz_data_all.empty:
            st.warning("No valid data available for visualizations after removing NaN values")
            return None
        
        # Extract clean data for visualization
        # Non-zero wells data
        well_names_non_zero = viz_data_non_zero[well_name_col]
        net_bo_data_non_zero = viz_data_non_zero[net_bo_col]
        net_diff_bo_data_non_zero = viz_data_non_zero[net_diff_bo_col]
        
        # All wells data
        well_names_all = viz_data_all[well_name_col]
        net_bo_data_all = viz_data_all[net_bo_col]
        
        # Check for finite values
        if (net_bo_data_non_zero.isna().all() or net_diff_bo_data_non_zero.isna().all() or 
            not np.isfinite(net_bo_data_non_zero).any() or not np.isfinite(net_diff_bo_data_non_zero).any() or
            net_bo_data_all.isna().all() or not np.isfinite(net_bo_data_all).any()):
            st.warning("No finite values available for visualization")
            return None
        
        # Create simplified subplots - 1 row, 3 columns for better layout
        fig, axes = plt.subplots(1, 3, figsize=(18, 6))
        fig.suptitle('Production Analysis Dashboard', fontsize=16, fontweight='bold')
        
        # 1. Net Diff BO by Well (Non-Zero Wells - Top 15)
        if len(net_diff_bo_data_non_zero) > 0 and len(well_names_non_zero) > 0:
            display_data = pd.DataFrame({
                'well_name': well_names_non_zero,
                'net_diff_bo': net_diff_bo_data_non_zero
            }).head(15)
            
            display_wells = display_data['well_name']
            display_net_diff = display_data['net_diff_bo']
            
            bars = axes[0].bar(range(len(display_wells)), display_net_diff, 
                              color=['lightgreen' if x >= 0 else 'lightcoral' for x in display_net_diff],
                              alpha=0.7)
            axes[0].set_xlabel('Wells')
            axes[0].set_ylabel('Net Diff BO')
            axes[0].set_title('Net Diff BO Performance (Top 15 Wells)')
            axes[0].set_xticks(range(len(display_wells)))
            axes[0].set_xticklabels(display_wells, rotation=45, ha='right')
            axes[0].grid(True, alpha=0.3)
            
            for bar, value in zip(bars, display_net_diff):
                height = bar.get_height()
                axes[0].text(bar.get_x() + bar.get_width()/2., height,
                            f'{value:.1f}', ha='center', va='bottom' if height >= 0 else 'top',
                            fontsize=8)
        else:
            axes[0].text(0.5, 0.5, 'No data available', ha='center', va='center', transform=axes[0].transAxes)
            axes[0].set_title('Net Diff BO Performance')
        
        # 2. Net BO by Well (Non-Zero Wells - Top 15)
        if len(net_bo_data_non_zero) > 0 and len(well_names_non_zero) > 0:
            display_data = pd.DataFrame({
                'well_name': well_names_non_zero,
                'net_bo': net_bo_data_non_zero
            }).head(15)
            
            display_wells = display_data['well_name']
            display_net_bo = display_data['net_bo']
            
            bars = axes[1].bar(range(len(display_wells)), display_net_bo, 
                              color='skyblue', alpha=0.7)
            axes[1].set_xlabel('Wells')
            axes[1].set_ylabel('Net BO')
            axes[1].set_title('Net BO Production (Top 15 Wells)')
            axes[1].set_xticks(range(len(display_wells)))
            axes[1].set_xticklabels(display_wells, rotation=45, ha='right')
            axes[1].grid(True, alpha=0.3)
            
            for bar, value in zip(bars, display_net_bo):
                height = bar.get_height()
                axes[1].text(bar.get_x() + bar.get_width()/2., height,
                            f'{value:.0f}', ha='center', va='bottom',
                            fontsize=8)
        else:
            axes[1].text(0.5, 0.5, 'No data available', ha='center', va='center', transform=axes[1].transAxes)
            axes[1].set_title('Net BO Production')
        
        # 3. Top 10 Wells with Highest Net BO (ALL WELLS)
        if len(net_bo_data_all) > 0 and len(well_names_all) > 0:
            # Get top 10 wells with highest Net BO from ALL wells
            top_wells_all = pd.DataFrame({
                'well_name': well_names_all,
                'net_bo': net_bo_data_all
            }).nlargest(10, 'net_bo')
            
            # Create horizontal bar chart for better readability
            bars = axes[2].barh(range(len(top_wells_all)), top_wells_all['net_bo'], 
                               color='gold', alpha=0.7, edgecolor='darkorange', linewidth=1)
            axes[2].set_xlabel('Net BO')
            axes[2].set_ylabel('Wells')
            axes[2].set_title('Top 10 Highest Producing Wells')
            axes[2].set_yticks(range(len(top_wells_all)))
            axes[2].set_yticklabels(top_wells_all['well_name'])
            axes[2].grid(True, alpha=0.3)
            
            # Add value labels on bars
            for bar, value in zip(bars, top_wells_all['net_bo']):
                width = bar.get_width()
                axes[2].text(width + width*0.01, bar.get_y() + bar.get_height()/2.,
                            f'{value:.0f}', ha='left', va='center', fontsize=9, fontweight='bold')
        else:
            axes[2].text(0.5, 0.5, 'No data available', ha='center', va='center', transform=axes[2].transAxes)
            axes[2].set_title('Top 10 Highest Producing Wells')
        
        plt.tight_layout()
        return fig
        
    except Exception as e:
        st.error(f"❌ Error creating visualizations: {str(e)}")
        import traceback
        st.error(f"Detailed error: {traceback.format_exc()}")
        return None

def create_comprehensive_powerpoint(data_df, well_count, stats, original_columns, visualization_fig):
    """
    Create a comprehensive PowerPoint presentation with data, statistics, and visualizations
    """
    try:
        # Create a new presentation
        prs = Presentation()
        
        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Production Analysis Report"
        subtitle.text = f"Comprehensive Well Performance Analysis\nTotal Wells: {stats['Total All Wells']}\nGenerated on: {pd.Timestamp.now().strftime('%Y-%m-%d %H:%M')}\nCreated by: Geol. Hassan Gamal Albery - Geologist @ Norpetco"
        
        # Executive Summary Slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Executive Summary"
        
        # Add summary content
        content_left = Inches(0.5)
        content_top = Inches(1.5)
        content_width = Inches(9.0)
        content_height = Inches(5.0)
        
        text_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        # Add summary points
        summary_points = [
            f"• Total Wells Analyzed: {stats['Total All Wells']}",
            f"• Wells with Non-Zero Net Diff BO: {stats['Total Wells with Non-Zero Net Diff BO']}",
            f"• Positive Performance Wells: {stats['Positive Net Diff BO Wells']}",
            f"• Wells Requiring Attention: {stats['Negative Net Diff BO Wells']}",
            f"• Total Net BO Production: {stats['Total Net BO (All Wells)']:,.0f}",
            f"• Average Net BO per Well: {stats['Average Net BO (All Wells)']:,.0f}",
            f"• Highest Producing Well: {stats['Maximum Net BO']:,.0f}",
            f"• Performance Range: {stats['Minimum Net BO']:,.0f} to {stats['Maximum Net BO']:,.0f}"
        ]
        
        for point in summary_points:
            p = text_frame.add_paragraph()
            p.text = point
            p.space_after = Inches(0.05)
        
        # Main Data Table Slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Production Data - Key Wells"
        
        # Create main data table (show only first 15 rows for readability, including TOTAL row if present)
        display_data = data_df.head(15) if len(data_df) > 15 else data_df
        
        rows = len(display_data) + 1
        cols = len(display_data.columns)
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9.0)
        height = Inches(0.8 * min(rows, 12))  # Limit height
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set column headers
        for i, column in enumerate(display_data.columns):
            table.cell(0, i).text = str(column)
        
        # Fill table with data
        for row_idx, (_, row_data) in enumerate(display_data.iterrows(), 1):
            for col_idx, column in enumerate(display_data.columns):
                value = row_data[column]
                if isinstance(value, (int, float)) and column not in [original_columns[0], original_columns[1]]:
                    table.cell(row_idx, col_idx).text = f"{value:,.2f}"
                else:
                    table.cell(row_idx, col_idx).text = str(value)
        
        # Key Metrics Slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Key Performance Metrics"
        
        # Create key metrics table
        key_metrics = {
            'Total Wells': stats['Total All Wells'],
            'Wells with Significant Changes': stats['Total Wells with Non-Zero Net Diff BO'],
            'Positive Performance Wells': stats['Positive Net Diff BO Wells'],
            'Wells Requiring Attention': stats['Negative Net Diff BO Wells'],
            'Total Net BO Production': stats['Total Net BO (All Wells)'],
            'Total Net Diff BO': stats['Total Net Diff BO (All Wells)'],
            'Average Net BO per Well': stats['Average Net BO (All Wells)'],
            'Highest Producing Well': stats['Maximum Net BO'],
            'Performance Standard Deviation': stats['Standard Deviation Net BO']
        }
        
        stats_rows = len(key_metrics) + 1
        stats_cols = 2
        left = Inches(1.0)
        top = Inches(1.5)
        width = Inches(8.0)
        height = Inches(0.8 * min(stats_rows, 15))
        
        stats_table = slide.shapes.add_table(stats_rows, stats_cols, left, top, width, height).table
        stats_table.cell(0, 0).text = "Metric"
        stats_table.cell(0, 1).text = "Value"
        
        for idx, (metric, value) in enumerate(key_metrics.items(), 1):
            stats_table.cell(idx, 0).text = metric
            if isinstance(value, (int, float)):
                if value > 1000:
                    stats_table.cell(idx, 1).text = f"{value:,.0f}"
                else:
                    stats_table.cell(idx, 1).text = f"{value:,.2f}"
            else:
                stats_table.cell(idx, 1).text = str(value)
        
        # Visualization Slides
        if visualization_fig:
            # Save figure to bytes
            img_buffer = io.BytesIO()
            visualization_fig.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
            img_buffer.seek(0)
            
            # Create individual visualization slides
            visualization_titles = [
                "Net Diff BO Performance",
                "Net BO Production", 
                "Top 10 Highest Producing Wells"
            ]
            
            for viz_title in visualization_titles:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = f"Analysis - {viz_title}"
                
                # Add the visualization image
                left = Inches(1.0)
                top = Inches(1.5)
                width = Inches(8.0)
                slide.shapes.add_picture(img_buffer, left, top, width=width)
        
        # Recommendations Slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Recommendations & Next Steps"
        
        text_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        recommendations = [
            "🎯 Focus Areas:",
            "• Analyze top performing wells for best practices replication",
            "• Review wells with negative Net Diff BO for improvement opportunities",
            "• Monitor wells with significant performance deviations",
            "",
            "📊 Operational Actions:",
            "• Optimize production parameters for underperforming wells",
            "• Implement preventive maintenance for critical wells",
            "• Share best practices from top performers",
            "",
            "📈 Continuous Improvement:",
            "• Regular monitoring of Net Diff BO trends",
            "• Periodic review of well performance categories",
            "• Update operational strategies based on performance data"
        ]
        
        for recommendation in recommendations:
            p = text_frame.add_paragraph()
            p.text = recommendation
            p.space_after = Inches(0.03)
        
        # Save to bytes buffer
        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        return ppt_buffer
        
    except Exception as e:
        st.error(f"❌ Error creating PowerPoint: {str(e)}")
        import traceback
        st.error(f"Detailed error: {traceback.format_exc()}")
        return None

def create_excel_with_visualizations(data_df, stats, visualization_fig):
    """
    Create an Excel file with data, statistics, and embedded visualizations
    """
    try:
        # Create Excel writer
        excel_buffer = io.BytesIO()
        
        with pd.ExcelWriter(excel_buffer, engine='xlsxwriter') as writer:
            # Write main data (include TOTAL row)
            data_df.to_excel(writer, sheet_name='Production Data', index=False)
            
            # Write statistics
            stats_df = pd.DataFrame(list(stats.items()), columns=['Metric', 'Value'])
            stats_df.to_excel(writer, sheet_name='Statistics', index=False)
            
            # Get workbook and worksheets
            workbook = writer.book
            
            # Format worksheets
            header_format = workbook.add_format({
                'bold': True,
                'text_wrap': True,
                'valign': 'top',
                'fg_color': '#D7E4BC',
                'border': 1
            })
            
            # Format data sheet
            data_sheet = writer.sheets['Production Data']
            for col_num, value in enumerate(data_df.columns.values):
                data_sheet.write(0, col_num, str(value), header_format)
            data_sheet.set_column('A:Z', 15)
            
            # Format statistics sheet
            stats_sheet = writer.sheets['Statistics']
            stats_sheet.write(0, 0, 'Metric', header_format)
            stats_sheet.write(0, 1, 'Value', header_format)
            stats_sheet.set_column('A:A', 35)
            stats_sheet.set_column('B:B', 20)
            
            # Add visualization if available
            if visualization_fig:
                # Save figure to bytes
                img_buffer = io.BytesIO()
                visualization_fig.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
                img_buffer.seek(0)
                
                # Create visualization sheet
                viz_sheet = workbook.add_worksheet('Visualizations')
                
                # Insert the image
                viz_sheet.insert_image('A1', 'visualization.png', {'image_data': img_buffer})
                viz_sheet.set_column('A:A', 50)
                viz_sheet.set_row(0, 300)
        
        excel_buffer.seek(0)
        return excel_buffer
        
    except Exception as e:
        st.error(f"❌ Error creating Excel file: {str(e)}")
        return None

def main():
    st.set_page_config(
        page_title="Production Analysis Dashboard", 
        page_icon="🛢️", 
        layout="wide",
        initial_sidebar_state="expanded"
    )
    
    # Enhanced Custom CSS for better user experience
    st.markdown("""
    <style>
    .main-header {
        font-size: 2.5rem;
        color: #1f77b4;
        text-align: center;
        margin-bottom: 2rem;
        font-weight: bold;
    }
    .info-box {
        background-color: #f8f9fa;
        padding: 1.5rem;
        border-radius: 10px;
        border-left: 5px solid #1f77b4;
        margin: 1rem 0;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
    }
    .metric-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 1.5rem;
        border-radius: 10px;
        box-shadow: 0 4px 6px rgba(0,0,0,0.1);
        text-align: center;
        color: white;
        margin: 0.5rem;
    }
    .upload-section {
        background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);
        padding: 2rem;
        border-radius: 15px;
        color: white;
        text-align: center;
        margin-bottom: 2rem;
    }
    .stButton button {
        width: 100%;
        border-radius: 8px;
        font-weight: bold;
        padding: 0.5rem 1rem;
    }
    .success-box {
        background-color: #d4edda;
        border: 1px solid #c3e6cb;
        border-radius: 8px;
        padding: 1rem;
        margin: 1rem 0;
    }
    .warning-box {
        background-color: #fff3cd;
        border: 1px solid #ffeaa7;
        border-radius: 8px;
        padding: 1rem;
        margin: 1rem 0;
    }
    .sidebar-developer {
        text-align: center;
        margin-bottom: 1rem;
        padding: 1rem;
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        border-radius: 10px;
        color: white;
    }
    </style>
    """, unsafe_allow_html=True)
    
    st.markdown('<h1 class="main-header">🛢️ Production Analysis Dashboard</h1>', unsafe_allow_html=True)
    
    # Enhanced Sidebar for navigation and info
    with st.sidebar:
        st.markdown("""
        <div style="text-align: center; margin-bottom: 1rem;">
            <span style="font-size: 3rem;">📊</span>
            <h2>Production Analytics</h2>
        </div>
        """, unsafe_allow_html=True)
        
        # Developer information under sidebar title
        st.markdown("""
        <div class="sidebar-developer">
        <h4>👨‍💻 Developed by</h4>
        <h3>Geol. Hassan Gamal Albery</h3>
        <p>Geologist @ Norpetco</p>
        </div>
        """, unsafe_allow_html=True)
        
        st.markdown("---")
        
        st.subheader("🚀 Quick Start")
        st.markdown("""
        1. **Upload** your Excel file
        2. **Review** automatic analysis
        3. **Download** reports
        """)
        
        st.markdown("---")
        st.subheader("📋 Supported Files")
        st.markdown("""
        • Excel (.xlsx)
        • Excel (.xls) 
        • Macro-enabled (.xlsm)
        """)
        
        st.markdown("---")
        st.subheader("🛠️ Tools")
        if st.button("🔄 Clear Cache & Refresh", use_container_width=True):
            st.runtime.legacy_caching.clear_cache()
            st.success("✅ Application refreshed!")
    
    # Main content area with improved layout
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.markdown("""
        <div class="info-box">
        <h3>🎯 What This Dashboard Does</h3>
        <p>This intelligent dashboard automatically analyzes your production data to provide:</p>
        <ul>
        <li><b>Well Performance Insights</b> - Identify top performers and areas for improvement</li>
        <li><b>Production Trends</b> - Track Net BO and Net Diff BO metrics</li>
        <li><b>Actionable Reports</b> - Download comprehensive analysis in multiple formats</li>
        <li><b>Visual Analytics</b> - Clear charts showing key performance indicators</li>
        </ul>
        </div>
        """, unsafe_allow_html=True)
    
    with col2:
        st.markdown("""
        <div class="info-box">
        <h3>📈 Key Features</h3>
        <p>• Automated Data Extraction<br>
           • Smart Column Detection<br>
           • Interactive Visualizations<br>
           • Multi-format Export<br>
           • Professional Reporting</p>
        </div>
        """, unsafe_allow_html=True)
    
    # Enhanced File upload section
    st.markdown("---")
    st.markdown('<div class="upload-section"><h2>📁 Upload Your Production Data</h2><p>Drag and drop your Excel file below to start analysis</p></div>', unsafe_allow_html=True)
    
    uploaded_file = st.file_uploader(
        "Choose your production Excel file", 
        type=['xlsx', 'xls', 'xlsm'],
        help="Upload an Excel file with production data. The app will automatically detect the required columns.",
        label_visibility="collapsed"
    )
    
    if uploaded_file is not None:
        try:
            # Process file without toggle status
            with st.spinner("🔄 Processing your file... This may take a few moments."):
                result_df, well_count, stats, original_columns, all_wells_data = extract_wells_with_net_diff_bo(uploaded_file)
                
                if result_df is not None and not result_df.empty:
                    # Generate visualizations (exclude TOTAL row for visualization)
                    data_without_total = result_df[result_df[original_columns[0]] != 'TOTAL (All Wells)']
                    fig = create_visualizations(data_without_total, original_columns, all_wells_data)
                    
                    # Generate PowerPoint automatically
                    ppt_buffer = create_comprehensive_powerpoint(result_df, well_count, stats, original_columns, fig)
                    
                    # Success message
                    st.markdown(f"""
                    <div class="success-box">
                    <h3>✅ Analysis Complete!</h3>
                    <p>Successfully processed <b>{stats['Total All Wells']}</b> total wells and identified <b>{well_count}</b> wells with significant Net Diff BO values.</p>
                    <p><b>PowerPoint report has been automatically generated and is ready for download below.</b></p>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    # Enhanced metrics display
                    st.markdown("---")
                    st.header("📊 Key Performance Indicators")
                    
                    kpi1, kpi2, kpi3, kpi4 = st.columns(4)
                    
                    with kpi1:
                        st.markdown(f"""
                        <div class="metric-card">
                        <h3>Total Wells</h3>
                        <h2>{stats['Total All Wells']}</h2>
                        </div>
                        """, unsafe_allow_html=True)
                    
                    with kpi2:
                        st.markdown(f"""
                        <div class="metric-card">
                        <h3>Wells with Changes</h3>
                        <h2>{stats['Total Wells with Non-Zero Net Diff BO']}</h2>
                        </div>
                        """, unsafe_allow_html=True)
                    
                    with kpi3:
                        st.markdown(f"""
                        <div class="metric-card">
                        <h3>Positive Performance</h3>
                        <h2>{stats['Positive Net Diff BO Wells']}</h2>
                        </div>
                        """, unsafe_allow_html=True)
                    
                    with kpi4:
                        st.markdown(f"""
                        <div class="metric-card">
                        <h3>Needs Attention</h3>
                        <h2>{stats['Negative Net Diff BO Wells']}</h2>
                        </div>
                        """, unsafe_allow_html=True)
                    
                    # Data preview (with TOTAL row included)
                    st.markdown("---")
                    st.header("📋 Production Data Overview")
                    st.dataframe(result_df, use_container_width=True, height=400)
                    
                    # Visualizations
                    st.markdown("---")
                    st.header("📈 Performance Analytics")
                    if fig:
                        st.pyplot(fig)
                        st.caption("Figure 1: Comprehensive production performance analysis across key metrics")
                    else:
                        st.info("📊 Visualizations not available due to insufficient data")
                    
                    # Enhanced Download section
                    st.markdown("---")
                    st.header("💾 Download Reports")
                    
                    st.markdown("""
                    <div class="info-box">
                    <h3>🎁 Export Your Analysis</h3>
                    <p>Choose from multiple formats to share your insights with your team:</p>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    download_col1, download_col2, download_col3 = st.columns(3)
                    
                    with download_col1:
                        st.subheader("📄 CSV Export")
                        st.markdown("Simple data format for spreadsheets")
                        # Export with TOTAL row included
                        csv = result_df.to_csv(index=False)
                        st.download_button(
                            label="📥 Download CSV",
                            data=csv,
                            file_name="production_analysis.csv",
                            mime="text/csv",
                            use_container_width=True
                        )
                    
                    with download_col2:
                        st.subheader("📊 Excel Report")
                        st.markdown("Complete analysis with charts")
                        if st.button("🔄 Generate Excel Report", use_container_width=True):
                            with st.spinner("Creating comprehensive Excel report..."):
                                excel_buffer = create_excel_with_visualizations(result_df, stats, fig)
                            
                            if excel_buffer:
                                st.download_button(
                                    label="📥 Download Excel",
                                    data=excel_buffer,
                                    file_name="production_analysis.xlsx",
                                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                    use_container_width=True
                                )
                            else:
                                st.error("❌ Failed to create Excel report")
                    
                    with download_col3:
                        st.subheader("🎤 PowerPoint")
                        st.markdown("Professional presentation")
                        if ppt_buffer:
                            st.download_button(
                                label="📥 Download PowerPoint",
                                data=ppt_buffer,
                                file_name="production_presentation.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                use_container_width=True
                            )
                            st.success("✅ PowerPoint ready for download!")
                        else:
                            st.error("❌ Failed to create PowerPoint presentation")
                
                else:
                    st.error("❌ No valid data found in the uploaded file. Please check your file format and try again.")
                    
        except Exception as e:
            st.error(f"❌ Error processing file: {str(e)}")
            st.markdown("""
            <div class="warning-box">
            <h3>💡 Troubleshooting Tips</h3>
            <ul>
            <li>Ensure your Excel file has data in the 'Report' worksheet</li>
            <li>Check that the file follows the expected format with multi-level headers</li>
            <li>Verify that required columns are present (Field, Well Names, Net BO, Net Diff BO)</li>
            <li>Try saving your file as .xlsx format if issues persist</li>
            </ul>
            </div>
            """, unsafe_allow_html=True)
    
    else:
        # Enhanced instructions when no file is uploaded
        st.markdown("---")
        st.header("📖 Getting Started Guide")
        
        guide_col1, guide_col2 = st.columns(2)
        
        with guide_col1:
            st.subheader("🎯 Step-by-Step Process")
            steps = [
                {"step": "1", "title": "Prepare Your Data", "desc": "Ensure your Excel file has production data in the 'Report' sheet with proper headers"},
                {"step": "2", "title": "Upload File", "desc": "Use the upload section above to select your Excel file (.xlsx, .xls, or .xlsm)"},
                {"step": "3", "title": "Automatic Analysis", "desc": "The app will automatically detect columns and process your data"},
                {"step": "4", "title": "Review Results", "desc": "Examine the insights, visualizations, and key metrics"},
                {"step": "5", "title": "Export Reports", "desc": "Download your analysis in CSV, Excel, or PowerPoint format"}
            ]
            
            for step in steps:
                with st.container():
                    st.markdown(f"**{step['step']}. {step['title']}**")
                    st.caption(step['desc'])
                    st.markdown("---")
        
        with guide_col2:
            st.subheader("📋 File Requirements")
            requirements = [
                "✅ **File Types**: .xlsx, .xls, or .xlsm (Macro-enabled Excel)",
                "✅ **Worksheet**: Data must be in 'Report' sheet",
                "✅ **Headers**: Multi-level headers (first 6 rows skipped)",
                "✅ **Required Columns**:",
                "   - Field information column",
                "   - Running wells names column", 
                "   - Net BO production values",
                "   - Net Diff BO performance values",
                "✅ **Data Format**: Stop processing at 'TOTAL' row"
            ]
            
            for req in requirements:
                st.markdown(req)
            
            st.markdown("---")
            st.subheader("🔍 Expected Output")
            st.markdown("""
            • **Data Table**: Filtered production data with TOTAL summary
            • **Key Metrics**: Performance statistics
            • **Visual Charts**: Three comprehensive visualizations
            • **Export Options**: Multiple report formats
            """)
        
        st.markdown("---")
        st.markdown("""
        <div style="text-align: center; padding: 2rem; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); border-radius: 15px; color: white;">
        <h2>🚀 Ready to Analyze Your Production Data?</h2>
        <p>Upload your Excel file above to unlock powerful insights and generate professional reports!</p>
        </div>
        """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()
